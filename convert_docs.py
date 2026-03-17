#!/usr/bin/env python3
"""
文档转换工具
把 PDF、Word、PPT 转成纯文本，方便 Claude 读取

使用方法：
1. 把文件放到 docs_to_convert/ 文件夹
2. 运行: python convert_docs.py
3. 转换后的文本在 converted_docs/ 文件夹
"""

import os
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation

# 输入输出目录
INPUT_DIR = Path("docs_to_convert")
OUTPUT_DIR = Path("converted_docs")

def convert_pdf(file_path):
    """转换 PDF 到文本"""
    text = []
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text())
        return '\n\n'.join(text)
    except Exception as e:
        return f"[转换失败: {e}]"

def convert_docx(file_path):
    """转换 Word 到文本"""
    try:
        doc = Document(file_path)
        text = [para.text for para in doc.paragraphs if para.text.strip()]
        return '\n\n'.join(text)
    except Exception as e:
        return f"[转换失败: {e}]"

def convert_pptx(file_path):
    """转换 PPT 到文本"""
    try:
        prs = Presentation(file_path)
        text = []
        for i, slide in enumerate(prs.slides, 1):
            text.append(f"--- 第 {i} 页 ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return '\n\n'.join(text)
    except Exception as e:
        return f"[转换失败: {e}]"

def main():
    # 创建目录
    INPUT_DIR.mkdir(exist_ok=True)
    OUTPUT_DIR.mkdir(exist_ok=True)

    # 扫描文件
    files = list(INPUT_DIR.glob("*"))
    if not files:
        print(f"📂 {INPUT_DIR}/ 文件夹是空的，请放入要转换的文件")
        return

    print(f"找到 {len(files)} 个文件，开始转换...\n")

    for file_path in files:
        if file_path.is_dir():
            continue

        suffix = file_path.suffix.lower()
        output_path = OUTPUT_DIR / f"{file_path.stem}.txt"

        print(f"📄 {file_path.name} ... ", end="")

        if suffix == '.pdf':
            content = convert_pdf(file_path)
        elif suffix in ['.docx', '.doc']:
            content = convert_docx(file_path)
        elif suffix in ['.pptx', '.ppt']:
            content = convert_pptx(file_path)
        else:
            print(f"⚠️  跳过（不支持的格式）")
            continue

        # 写入文件
        output_path.write_text(content, encoding='utf-8')
        print(f"✅ → {output_path.name}")

    print(f"\n✨ 完成！转换后的文件在 {OUTPUT_DIR}/ 文件夹")

if __name__ == "__main__":
    main()
